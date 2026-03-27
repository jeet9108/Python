import sys
from pptx import Presentation

def extract_text(file_paths, out_path):
    with open(out_path, "w", encoding="utf-8") as f:
        for file_path in file_paths:
            f.write(f"--- {file_path} ---\n")
            try:
                prs = Presentation(file_path)
                for i, slide in enumerate(prs.slides):
                    f.write(f"Slide {i+1}:\n")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            f.write("  " + shape.text.replace("\n", " | ") + "\n")
            except Exception as e:
                f.write(f"Error reading {file_path}: {e}\n")

if __name__ == "__main__":
    extract_text(sys.argv[1:-1], sys.argv[-1])
